from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
BLUE = (66, 133, 244)
GREEN = (52, 168, 83)
YELLOW = (251, 188, 4)
RED = (234, 67, 53)
DARK = (16, 20, 24)
PANEL = (27, 32, 39)
TEXT = (248, 250, 252)
MUTED = (203, 213, 225)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def draw_box(draw: ImageDraw.ImageDraw, xy: tuple[int, int, int, int], title: str, body: str, color: tuple[int, int, int]) -> None:
    draw.rounded_rectangle(xy, radius=18, fill=PANEL, outline=color, width=4)
    x1, y1, x2, _ = xy
    draw.ellipse((x1 + 24, y1 + 24, x1 + 64, y1 + 64), fill=color)
    draw.text((x1 + 82, y1 + 20), title, fill=TEXT, font=font(30, True))
    words = body.split()
    lines, current = [], ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if draw.textlength(candidate, font=font(21)) < (x2 - x1 - 54):
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    for idx, line in enumerate(lines[:3]):
        draw.text((x1 + 28, y1 + 78 + idx * 30), line, fill=MUTED, font=font(21))


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: tuple[int, int, int]) -> None:
    draw.line((start, end), fill=color, width=5)
    ex, ey = end
    draw.polygon([(ex, ey), (ex - 18, ey - 10), (ex - 18, ey + 10)], fill=color)


def make_architecture() -> None:
    img = Image.new("RGB", (1800, 1100), DARK)
    draw = ImageDraw.Draw(img)
    draw.text((70, 54), "AccessPrep AI Architecture", fill=TEXT, font=font(56, True))
    draw.text((74, 126), "Voice-first decision intelligence for accessible competitive exam preparation", fill=MUTED, font=font(28))

    boxes = [
        ((70, 240, 390, 430), "Student", "Screen reader, keyboard, audio output, mobile or desktop", BLUE),
        ((500, 240, 840, 430), "Streamlit UI", "Home, Upload, AI Tutor, Study Recommendation, About", GREEN),
        ((950, 160, 1310, 340), "OCR Layer", "Google Vision OCR or Gemini Vision extracts text", YELLOW),
        ((950, 420, 1310, 600), "Gemini Tutor", "Explain, summarize, answer follow-ups, generate MCQs", BLUE),
        ((1420, 290, 1730, 500), "Decision Engine", "Weak areas, next chapter, revision plan, confidence", RED),
        ((500, 620, 840, 810), "Accessibility", "gTTS speech, large fonts, dark mode, voice-friendly output", GREEN),
        ((950, 710, 1310, 900), "Cloud Run", "Containerized Python app with secure environment variables", BLUE),
    ]
    for box in boxes:
        draw_box(draw, *box)

    draw_arrow(draw, (390, 335), (500, 335), BLUE)
    draw_arrow(draw, (840, 310), (950, 250), YELLOW)
    draw_arrow(draw, (840, 360), (950, 500), BLUE)
    draw_arrow(draw, (1310, 500), (1420, 395), RED)
    draw_arrow(draw, (670, 430), (670, 620), GREEN)
    draw_arrow(draw, (1125, 600), (1125, 710), BLUE)
    draw.text((70, 1000), "Google Cloud style deployment: Cloud Run + Secret Manager + Vision API + Gemini API", fill=MUTED, font=font(25))
    img.save(ROOT / "architecture.png")


def set_text(shape, text: str, size: int = 24, bold: bool = False, color: tuple[int, int, int] = TEXT) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    frame = shape.text_frame
    frame.clear()
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.7))
    set_text(box, title, 34, True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.05), Inches(11.8), Inches(0.45))
        set_text(sub, subtitle, 16, False, MUTED)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, color: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*PANEL)
    shape.line.color.rgb = RGBColor(*color)
    shape.line.width = Pt(1.5)
    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.2), Inches(w - 0.44), Inches(0.35))
    set_text(title_box, title, 18, True, color)
    body_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.68), Inches(w - 0.44), Inches(h - 0.82))
    set_text(body_box, body, 15, False, TEXT)


def make_ppt() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    out_dir = ROOT / "ppt"
    out_dir.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*DARK)
        return slide

    slide = blank()
    add_title(slide, "AccessPrep AI", "AI-powered Decision Intelligence for Accessible Competitive Exam Preparation")
    add_card(slide, 0.8, 2.0, 5.8, 2.8, "Inclusive Intelligence Lab", "AI for Every Ability.\n\nA multidisciplinary team building accessible, inclusive, human-centered AI technologies.", GREEN)
    add_card(slide, 7.0, 2.0, 5.2, 2.8, "Hackathon Demo", "Computer Vision + OCR + Gemini + Conversational AI + Voice-first accessibility + Personalized study recommendations.", BLUE)

    slides = [
        ("Problem", [("Accessibility Gap", "Scanned PDFs and image notes are difficult for visually impaired aspirants to use independently.", RED), ("Learning Gap", "OCR alone does not explain concepts, answer questions, create quizzes, or plan revision.", YELLOW)]),
        ("Solution", [("Unified Workflow", "Upload material, extract text, learn with Gemini, generate MCQs, and listen to responses.", BLUE), ("Decision Intelligence", "Profile-aware recommendations identify weak areas, next chapter, revision topics, and confidence.", GREEN)]),
        ("Architecture", [("Pipeline", "Student interface -> Streamlit -> OCR layer -> Gemini tutor -> decision engine -> speech output.", BLUE), ("Cloud Ready", "Dockerized Python application deployable on Google Cloud Run with API keys as secrets.", GREEN)]),
        ("Technology Stack", [("Google Cloud", "Cloud Run, Vertex AI Gemini, Gemini API, Vision API, BigQuery, Artifact Registry, Secret Manager, Cloud Logging.", BLUE), ("Application", "Python, Streamlit, python-dotenv, gTTS, Pillow, PyMuPDF, Docker.", YELLOW)]),
        ("Demo", [("1. Upload", "Upload an image or PDF and extract accessible text.", GREEN), ("2. Tutor", "Explain, summarize, ask follow-up questions, and generate five MCQs.", BLUE), ("3. Recommend", "Generate weak topics, next chapter, revision schedule, study plan, and confidence.", RED)]),
        ("Impact", [("Independent Learning", "Students can convert inaccessible material into audio-friendly explanations and practice.", GREEN), ("Exam Readiness", "Personalized study decisions help focus limited preparation time on high-value topics.", BLUE)]),
        ("Future Scope", [("Inclusive Expansion", "Speech input, regional languages, Braille display support, offline study packs.", YELLOW), ("Institutional Scale", "Teacher dashboards, adaptive mock tests, progress analytics, and community mentoring.", GREEN)]),
    ]
    for idx, (title, cards) in enumerate(slides, start=2):
        slide = blank()
        add_title(slide, f"Slide {idx}: {title}")
        if title == "Architecture":
            slide.shapes.add_picture(str(ROOT / "architecture.png"), Inches(0.65), Inches(1.55), width=Inches(12.0))
        else:
            for card_index, (card_title, body, color) in enumerate(cards):
                add_card(slide, 0.8 + (card_index % 3) * 4.1, 2.0, 3.75, 2.55, card_title, body, color)

    prs.save(out_dir / "presentation.pptx")


if __name__ == "__main__":
    make_architecture()
    make_ppt()
